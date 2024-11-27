# Import necessary libraries
import streamlit as st
from langchain.document_loaders import PyPDFLoader, TextLoader
from langchain.vectorstores import FAISS
from langchain.schema import Document
from langchain.chains import RetrievalQA
from langchain.llms import OpenAI
from pptx import Presentation
import pandas as pd
from sentence_transformers import SentenceTransformer
from transformers import pipeline
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate
from langchain.chains.question_answering import load_qa_chain
from langchain.llms import Ollama
import os
import tempfile
from collections import Counter
import json

# <-----------------------------------Set up Streamlit app------------------------------------>
st.set_page_config(page_title="Corporate Training Knowledge Hub", layout="wide")
st.title("Corporate Training Knowledge Hub")

# <------------------------------------Initialize components------------------------------------->
llm = Ollama(model="llama3.2")  # Replace with your Llama model
embedding_model = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
vectorstore = None
document_store = []

# <---------------------------------------------Define tabs for functionalities------------------------------------>
tabs = st.tabs(["Upload Files", "Original Context", "Document Summarization", "Interactive Q&A", "Word Cloud"])

# <--------------------------------------------------Upload and process files------------------------------------->

def process_files(uploaded_files):
    global vectorstore
    for uploaded_file in uploaded_files:
        file_type = uploaded_file.name.split(".")[-1]
        documents = []

        # Save the uploaded file to a temporary location
        with tempfile.NamedTemporaryFile(delete=False) as temp_file:
            temp_file.write(uploaded_file.read())
            temp_file_path = temp_file.name

        if file_type == "pdf":
            loader = PyPDFLoader(temp_file_path)
            documents = loader.load()

        elif file_type == "pptx":
            presentation = Presentation(temp_file_path)
            content = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        content += shape.text_frame.text + " "
            documents = [Document(page_content=content, metadata={"name": uploaded_file.name})]

        elif file_type == "txt":
            with open(temp_file_path, "r", encoding="utf-8") as file:
                content = file.read()
            documents = [Document(page_content=content, metadata={"name": uploaded_file.name})]

        elif file_type == "xlsx":
            excel_data = pd.read_excel(temp_file_path)
            content = excel_data.to_string(index=False)
            documents = [Document(page_content=content, metadata={"name": uploaded_file.name})]

        # If metadata is missing, set a default name
        for doc in documents:
            if "metadata" not in doc or "name" not in doc.metadata:
                doc.metadata = doc.metadata or {}
                doc.metadata["name"] = uploaded_file.name

        # Store documents and update FAISS
        document_store.extend(documents)
        texts = [doc.page_content for doc in documents]
        if vectorstore is None:
            vectorstore = FAISS.from_texts(texts, embedding_model)
        else:
            vectorstore.add_texts(texts)

        # Delete the temporary file after processing
        os.remove(temp_file_path)

# <----------------------------------------------------Summarization function------------------------------------->

summary_prompt = PromptTemplate(
    input_variables=["context"],
    template=(
        "You are a corporate training assistant. Summarize the following text into clear, concise, and "
        "professional 5 to 8 sentences. Focus on key details, actionable insights, and important takeaways relevant "
        "to corporate training topics. Maintain a formal tone:\n\n{context}"
    ),
)

def summarize_text(text):
    if not vectorstore:
        return "No documents indexed for summarization. Please upload files first."

    retriever = vectorstore.as_retriever()

    # Retrieve relevant documents
    docs = retriever.get_relevant_documents(text)
    if not docs:
        return "No relevant documents found for summarization."

    # Combine documents into a single context
    combined_context = " ".join([doc.page_content for doc in docs])

    # Chunk the combined context if it's too large for the LLM to process
    max_chunk_size = 2048
    chunks = [combined_context[i:i + max_chunk_size] for i in range(0, len(combined_context), max_chunk_size)]

    # Create an LLM chain for summarization
    llm_chain = LLMChain(llm=llm, prompt=summary_prompt)

    # Summarize each chunk and combine the results
    summaries = []
    for chunk in chunks:
        try:
            response = llm_chain.run({"context": chunk})
            summaries.append(response.strip())
        except Exception as e:
            summaries.append(f"Error processing chunk: {str(e)}")

    final_summary = " ".join(summaries)
    return final_summary

# <------------------------------------------------------Q&A function----------------------------------->
def answer_question(question):
    if not document_store:
        return "No documents uploaded yet. Please upload files first."
    if vectorstore is None:
        return "No documents indexed for retrieval. Please upload files."

    retriever = vectorstore.as_retriever()
    qa_chain = RetrievalQA(llm=llm, retriever=retriever)
    return qa_chain.run(question)

# <-------------------------------------------- Word Cloud Function----------------------------------->

def get_word_frequencies(text):
    words = text.split()
    word_counts = Counter(words)
    return [{"word": word, "count": count} for word, count in word_counts.items()]

# <-------------------------------------------------------Main App-------------------------------->

st.sidebar.header("Welcome!")
st.sidebar.info(
    "Upload corporate training documents, explore their contents, get concise summaries, "
    "generate word clouds, and interactively query the knowledge base!"
)

with tabs[0]:
    st.header("Upload Files")
    uploaded_files = st.file_uploader(
        "Upload corporate documents (PDF, PPTX, TXT, XLSX)", type=["pdf", "pptx", "txt", "xlsx"], accept_multiple_files=True
    )
    if uploaded_files:
        process_files(uploaded_files)
        st.success("Files processed successfully!")

with tabs[1]:
    st.header("Original Context")
    if document_store:
        for doc in document_store:
            # Safely access document metadata
            doc_name = doc.metadata.get("name", "Unknown Document")
            st.write(f"### {doc_name}")
            st.text_area("Content", doc.page_content, height=300)
    else:
        st.info("Please upload files to display their content.")

with tabs[2]:
    st.header("Document Summarization")
    if document_store:
        text_data = " ".join([doc.page_content for doc in document_store])
        summary = summarize_text(text_data)
        st.write("### Summary:")
        st.write(summary)
    else:
        st.info("Please upload files to summarize.")

with tabs[3]:
    st.header("Interactive Q&A")
    question = st.text_input("Ask a question about the uploaded documents:")
    if question:
        answer = answer_question(question)
        st.write("### Answer:")
        st.write(answer)

with tabs[4]:
    st.header("Dynamic Word Cloud")
    uploaded_file = st.file_uploader("Upload a text file for Word Cloud Animation", type=["txt"])
    if uploaded_file:
        text = uploaded_file.read().decode("utf-8")
        word_data = get_word_frequencies(text)
        word_data_json = json.dumps(word_data)

        st.markdown("### Animated Word Cloud")
        st.components.v1.html(f"""
        <!DOCTYPE html>
        <html>
        <head>
            <script src="https://cdnjs.cloudflare.com/ajax/libs/p5.js/1.6.0/p5.min.js"></script>
            <style>
                body {{ margin: 0; overflow: hidden; }}
            </style>
        </head>
        <body>
            <script>
                const wordData = {word_data_json};
                let bubbles = [];

                function setup() {{
                    createCanvas(windowWidth, windowHeight);
                    wordData.forEach(data => {{
                        let size = map(data.count, 1, Math.max(...wordData.map(w => w.count)), 20, 100);
                        bubbles.push(new Bubble(data.word, random(width), height + size, size));
                    }});
                }}

                function draw() {{
                    background(255);
                    bubbles.forEach(bubble => {{
                        bubble.move();
                        bubble.display();
                    }});
                }}

                class Bubble {{
                    constructor(word, x, y, size) {{
                        this.word = word;
                        this.x = x;
                        this.y = y;
                        this.size = size;
                        this.speed = random(1, 3);
                    }}

                    move() {{
                        this.y -= this.speed;
                        if (this.y < -this.size) {{
                            this.y = height + this.size;
                            this.x = random(width);
                        }}
                    }}

                    display() {{
                        fill(100, 100, 255, 150);
                        noStroke();
                        ellipse(this.x, this.y, this.size);
                        fill(0);
                        textAlign(CENTER, CENTER);
                        textSize(this.size / 4);
                        text(this.word, this.x, this.y);
                    }}
                }}
            </script>
        </body>
        </html>
        """, height=600)
    else:
        st.info("Please upload a text file to generate the animated Word Cloud.")

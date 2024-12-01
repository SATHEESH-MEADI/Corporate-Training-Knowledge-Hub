
# Import necessary libraries
import streamlit as st
from langchain_community.document_loaders import PyPDFLoader, TextLoader  # Updated imports
from langchain_community.vectorstores import FAISS  # Updated imports
from langchain.schema import Document
from langchain.chains import RetrievalQA
from langchain_community.llms import OpenAI, Ollama  # Updated imports
from pptx import Presentation
import pandas as pd
from sentence_transformers import SentenceTransformer
from transformers import pipeline
from langchain_community.embeddings import HuggingFaceEmbeddings  # Updated imports
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate
from langchain.chains.question_answering import load_qa_chain
from langchain_community.llms import HuggingFacePipeline  # Updated imports
import os
from wordcloud import WordCloud
import matplotlib.pyplot as plt
from transformers import BartForConditionalGeneration, BartTokenizer
from difflib import HtmlDiff, SequenceMatcher
import tempfile
from transformers import AutoModelForCausalLM, AutoTokenizer, AutoModelForTokenClassification, AutoModelForSeq2SeqLM
from transformers import pipeline
import nltk
from langchain_core.messages import HumanMessage,AIMessage
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.prompts import MessagesPlaceholder
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain.chains import create_retrieval_chain
from langchain.chains import create_history_aware_retriever
from langchain.memory import ConversationBufferMemory
import random
import re
import pprint
import json









# Download NLTK punkt tokenizer for sentence splitting
nltk.download("punkt")
from nltk.tokenize import sent_tokenize


# <-----------------------------------Set up Streamlit app------------------------------------>
st.set_page_config(page_title="Corporate Training Knowledge Hub", layout="wide")
st.title("Corporate Training Knowledge Hub")

# <------------------------------------Initialize components------------------------------------->
# Initialize the LLaMA model using Ollama
llm = Ollama(model="llama3.2")  # Replace with your locally installed LLaMA model


    


# Load embedding model for document processing
embedding_model = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
vectorstore = None
document_store = []

# <---------------------------------------------Define tabs for functionalities------------------------------------>

tabs = st.tabs(["Upload Files", "Original Context", "Document Summarization", "Interactive Chatbot", "Quiz", "Compare Docs","Word Cloud", "Highlights"])

# <--------------------------------------------------Upload and process files------------------------------------->
def process_files(uploaded_files):
    global vectorstore
    for uploaded_file in uploaded_files:
        file_type = uploaded_file.name.split(".")[-1]
        combined_content = ""
        document = None

        with tempfile.NamedTemporaryFile(delete=False) as temp_file:
            temp_file.write(uploaded_file.read())
            temp_file_path = temp_file.name

        if file_type == "pdf":
            loader = PyPDFLoader(temp_file_path)
            all_pages = loader.load()
            combined_content = " ".join([page.page_content for page in all_pages])
            document = Document(page_content=combined_content, metadata={"name": uploaded_file.name})

        elif file_type == "pptx":
            presentation = Presentation(temp_file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        combined_content += shape.text_frame.text + " "
            document = Document(page_content=combined_content, metadata={"name": uploaded_file.name})

        elif file_type == "txt":
            with open(temp_file_path, "r", encoding="utf-8") as file:
                combined_content = file.read()
            document = Document(page_content=combined_content, metadata={"name": uploaded_file.name})

        elif file_type == "xlsx":
            excel_data = pd.read_excel(temp_file_path)
            combined_content = excel_data.to_string(index=False)
            document = Document(page_content=combined_content, metadata={"name": uploaded_file.name})

        if document is None:
            st.warning(f"File type '{file_type}' is not supported.")
            os.remove(temp_file_path)
            continue

        document_store.append(document)
        texts = [document.page_content]
        if vectorstore is None:
            vectorstore = FAISS.from_texts(texts, embedding_model)
        else:
            vectorstore.add_texts(texts)

        os.remove(temp_file_path)

# <----------------------------------------------------Summarization function------------------------------------->

def summarize_text_with_llama(text):
    """
    Summarizes the provided text using the locally running LLaMA 3.2 model.
    """
    # Prepare the context for summarization
    prompt = f"""
    Please summarize the following text:
    
    {text}
    """

    # Use the locally running LLaMA model to generate the summary
    response = llm(prompt)

    # Extract and return the generated summary
    summary = response.strip()
    return summary

# <------------------------------------------------------Interactive Q&A Functionality----------------------------------->

def answer_question_with_llama(question):
    """
    Answers a question using the locally installed LLaMA model with Ollama.
    """
    # Guard clause for empty/None question
    if not question:
        return "Please provide a question to answer."
    
    # Guard clause for vectorstore
    if not vectorstore:
        return "No documents indexed for retrieval. Please upload files first."

    try:
        # Retrieve relevant documents
        retriever = vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": 3})
        retrieved_docs = retriever.get_relevant_documents(question)

        if not retrieved_docs:
            return "No relevant documents found for your question."
        
        prompt = ChatPromptTemplate.from_messages([
            MessagesPlaceholder(variable_name="chat_history"),
            ("user", "{input}"),
            ("user", "Given the above conversation, generate a search query to look up in order to get information relevant to the conversation")
        ])

        history_retriever_chain = create_history_aware_retriever(llm, retriever, prompt)

        answer_prompt = ChatPromptTemplate.from_messages([
            ("system", "Answer the user's questions based on the below context:\n\n{context}"),
            MessagesPlaceholder(variable_name="chat_history"),
            ("user", "{input}")
        ])

        # Create the document processing chain
        document_chain = create_stuff_documents_chain(llm, answer_prompt)

        # Create the final conversational retrieval chain
        conversational_retrieval_chain = create_retrieval_chain(
            history_retriever_chain, 
            document_chain
        )

        # Initialize chat history if not exists
        chat_history = []

        # Invoke the chain with the question
        response = conversational_retrieval_chain.invoke({
            'chat_history': chat_history,
            "input": question
        })

        # Update chat history
        chat_history.append((
            HumanMessage(content=question), 
            AIMessage(content=response["answer"])
        ))

        return response['answer']

    except Exception as e:
        return f"An error occurred while processing your question: {str(e)}"


# <-------------------------------------------- Word Cloud Function---------------------------------->
def generate_word_cloud(text):
    """
    Generates a word cloud visualization for the provided text.
    """
    wordcloud = WordCloud(background_color="white", width=800, height=400, max_words=200, colormap="viridis").generate(text)
    plt.figure(figsize=(10, 5))
    plt.imshow(wordcloud, interpolation="bilinear")
    plt.axis("off")
    st.pyplot(plt)

# <------------------------------------------------------------Compare Documents------------------------------------------------------>
def compare_documents():
    """
    Compares the first two documents in the document store for similarities and differences.
    """
    if len(document_store) < 2:
        st.error("Please upload at least two documents to compare.")
        return

    doc1_content = document_store[0].page_content
    doc2_content = document_store[1].page_content

    differ = HtmlDiff()
    html_diff = differ.make_file(doc1_content.splitlines(), doc2_content.splitlines(), 
                                 fromdesc=document_store[0].metadata['name'], todesc=document_store[1].metadata['name'])
    st.write("### Comparison Result")
    st.components.v1.html(html_diff, height=600, scrolling=True)
    similarity = SequenceMatcher(None, doc1_content, doc2_content).ratio()
    st.write(f"### Similarity Score: {similarity:.2%}")


#< -----------------------------------------------------------Highlights-------------------------------------->


ner_model = pipeline("ner", model="dbmdz/bert-large-cased-finetuned-conll03-english", aggregation_strategy="simple")


def generate_description_with_ollama(entity_text, entity_type, context):
    """
    Generates concise descriptions for entities using Ollama.
    """
    prompt = f"""
    Entity: {entity_text} ({entity_type})
    Context: {context}
    Task: Provide a single-sentence refined description of this entity.
    """

    # Use the locally installed Ollama model
    response = llm(prompt)
    return response.strip()

def extract_highlights_with_ollama(text):
    """
    Extracts concise highlights and generates descriptions using Ollama 3.2.
    """
    entities = ner_model(text)
    valid_entity_types = {"PER", "ORG", "LOC", "GPE", "DATE"}
    seen_entities = set()
    highlights = []

    for entity in entities:
        entity_text = entity["word"]
        entity_type = entity["entity_group"]

        if entity_type in valid_entity_types and entity_text not in seen_entities:
            seen_entities.add(entity_text)

            # Extract context sentences (1-2 sentences only)
            context_sentences = [s for s in sent_tokenize(text) if entity_text in s]
            context = " ".join(context_sentences[:1]) if context_sentences else "No detailed context available."

            # Generate concise description
            refined_description = generate_description_with_ollama(entity_text, entity_type, context)
            highlights.append(f"{entity_text} ({entity_type}) - {refined_description}")

    return highlights

#<---------------------------------------------------Quiz---------------------------------------->
import re
import streamlit as st


# Function to parse quiz questions
def parse_quiz(response):
    quiz_questions = []
    try:
        # Split the response into individual questions using regex
        questions = re.split(r"(?=Question \d+:)", response)

        for question in questions:
            if not question.strip():
                continue

            # Extract the question text
            question_match = re.search(r"Question \d+: (.+)", question)
            question_text = question_match.group(1).strip() if question_match else None

            # Extract choices
            choices = {}
            choice_matches = re.findall(r"([A-D])\) (.+)", question)

            for choice in choice_matches:
                choices[choice[0]] = choice[1].strip()

            # Extract the answer
            answer_match = re.search(r"Answer: ([A-D])", question)
            correct_answer = answer_match.group(1).strip() if answer_match else None

            if question_text and choices and correct_answer:
                quiz_questions.append({
                    "question": question_text,
                    "choices": [f"{key}) {value}" for key, value in choices.items()],
                    "answer": correct_answer
                })

        return quiz_questions

    except Exception as e:
        raise ValueError(f"Error while parsing quiz questions: {e}")
def generate_quiz_questions(document_content, num_questions=5):
    """
    Generate quiz questions based on the document content.
    """
    # Construct the prompt for the LLM
    question_prompt = (
        f"Generate {num_questions} multiple-choice quiz questions from the following text:\n\n{document_content}\n\n"
        "Format the output as:\n"
        "Question [number]: [Your question here]\n"
        "Choices:\n"
        "A) [Option A]\n"
        "B) [Option B]\n"
        "C) [Option C]\n"
        "D) [Option D]\n"
        "Answer: [Correct option letter]\n"
        "Ensure all questions are relevant and based on the text provided."
    )

    # Get the response from the LLM
    response = llm(question_prompt)

    # Validate and parse the structured response
    try:
        quiz_questions = parse_quiz(response)
        return quiz_questions
    except Exception as e:
        st.error(f"Error parsing quiz questions: {e}")
        return []


def display_quiz_with_checkboxes(quiz_questions):
    """
    Display all questions using checkboxes for answer selection within a form.
    """
    # Initialize session state for answers and feedback
    if "user_answers" not in st.session_state:
        st.session_state.user_answers = {idx: [] for idx in range(len(quiz_questions))}
    if "submitted" not in st.session_state:
        st.session_state.submitted = False

    st.title("Quiz")

    # Start a form to batch interactions
    with st.form("quiz_form"):
        for idx, question_data in enumerate(quiz_questions):
            st.write(f"### Question {idx + 1}: {question_data['question']}")

            selected_choices = st.session_state.user_answers.get(idx, [])

            # Display choices as checkboxes
            updated_choices = []
            for choice in question_data["choices"]:
                is_checked = choice in selected_choices
                if st.checkbox(label=choice, value=is_checked, key=f"checkbox_q{idx}_{choice}"):
                    updated_choices.append(choice)

            # Update session state for selected answers
            st.session_state.user_answers[idx] = updated_choices

        # Submit button to finalize answers
        submitted = st.form_submit_button("Submit Quiz")
        if submitted:
            st.session_state.submitted = True

    # Display feedback after submission
    if st.session_state.submitted:
        st.write("## Quiz Summary:")
        for idx, question_data in enumerate(quiz_questions):
            user_answer = st.session_state.user_answers[idx]
            correct_answer = question_data["answer"]
            correct_option = f"{correct_answer}) " + next(
                c.split(") ", 1)[1] for c in question_data["choices"] if c.startswith(correct_answer)
            )
            is_correct = "✔️" if correct_option in user_answer and len(user_answer) == 1 else "❌"
            st.write(
                f"**Q{idx + 1}: {question_data['question']}**\n"
                f"Your Answer: {', '.join(user_answer) or 'No Answer'}\n"
                f"Correct Answer: {correct_option} {is_correct}\n"
            )





# <-------------------------------------------------------Main App-------------------------------->
st.sidebar.header("Welcome!")
st.sidebar.info("Upload corporate training documents, explore their contents, get concise summaries, generate word clouds, and ask interactive questions!")

with tabs[0]:
    st.header("Upload Files")
    uploaded_files = st.file_uploader("Upload corporate documents (PDF, PPTX, TXT, XLSX)", 
                                      type=["pdf", "pptx", "txt", "xlsx"], accept_multiple_files=True)
    if uploaded_files:
        process_files(uploaded_files)
        st.success("Files processed successfully!")

with tabs[1]:
    st.header("Original Context")
    if document_store:
        for doc in document_store:
            st.write(f"### {doc.metadata['name']}")
            st.text_area(f"Content_{doc.metadata['name']}", doc.page_content, height=300)
    else:
        st.info("Please upload files to display their content.")

with tabs[2]:
    st.header("Document Summarization")
    if document_store:
        for doc in document_store:
            st.write(f"### {doc.metadata['name']}")
            summary = summarize_text_with_llama(doc.page_content)
            st.write("### Summary:")
            st.write(summary)
    else:
        st.info("Please upload files to summarize.")

with tabs[3]:
    st.header("Interactive Q&A")
    
    # Check for vectorstore
    if not vectorstore:
        st.info("Please upload documents to enable the Q&A functionality.")
        st.stop()

    # Initialize session state for messages if not exists
    if 'messages' not in st.session_state:
        st.session_state.messages = []

    # Display existing messages
    for message in st.session_state.messages:
        with st.chat_message(message['role']):
            st.markdown(message['content'])

    # Handle new user input
    if prompt := st.chat_input('Ask questions about the uploaded document(s)'):
        # Fixed typo in session_state
        st.session_state.messages.append({'role': '👽', 'content': prompt})

        with st.chat_message('user'):
            st.markdown(prompt)

        with st.chat_message('assistant'):
            with st.spinner('Generating response...'):
                result = answer_question_with_llama(prompt)
                st.write(result)

        st.session_state.messages.append({'role': '🤖', 'content': result})


with tabs[6]:
    st.header("Word Cloud")
    if document_store:
        text_data = " ".join([doc.page_content for doc in document_store])
        st.write("### Word Cloud")
        generate_word_cloud(text_data)
    else:
        st.info("Please upload files to generate a word cloud.")

with tabs[5]:
    st.header("Compare Documents")
    compare_documents()

with tabs[7]:
    st.header("Highlights (Concise Contextual Insights)")

    if document_store:
        for doc in document_store:
            st.subheader(f"Document: {doc.metadata['name']}")
            highlights = extract_highlights_with_ollama(doc.page_content)
            if highlights:
                for i, highlight in enumerate(highlights, start=1):
                    st.markdown(f"**{i}. {highlight}**")
            else:
                st.write("No significant entities or highlights found.")
    else:
        st.info("No documents uploaded yet.")

# Integration with Streamlit Tabs
with tabs[4]:
    st.header("Quiz")
    if document_store:
        # Allow the user to select a document for the quiz
        quiz_document = st.selectbox("Select a document for the quiz:", [doc.metadata["name"] for doc in document_store])
        selected_doc = next(doc for doc in document_store if doc.metadata["name"] == quiz_document)

        # Check if questions have already been generated for the selected document
        if "quiz_questions" not in st.session_state:
            st.session_state.quiz_questions = generate_quiz_questions(selected_doc.page_content)

        # Display the stored questions
        quiz_questions = st.session_state.quiz_questions

        # Display the quiz with checkboxes
        if quiz_questions:
            display_quiz_with_checkboxes(quiz_questions)
        else:
            st.warning("No quiz questions were generated. Please check the document content.")
    else:
        st.info("Please upload documents to create quizzes.")

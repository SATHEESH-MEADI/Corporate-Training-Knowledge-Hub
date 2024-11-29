# Import necessary libraries
import streamlit as st
from langchain_community.document_loaders import PyPDFLoader, TextLoader  # Updated imports
from langchain_community.vectorstores import FAISS  # Updated imports
from langchain.schema import Document
from langchain.chains import RetrievalQA
from langchain_community.llms import OpenAI, Ollama  # Updated imports
from pptx import Presentation
import pandas as pd
from sentence_transformers import SentenceTransformer
from transformers import pipeline
from langchain_community.embeddings import HuggingFaceEmbeddings  # Updated imports
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate
from langchain.chains.question_answering import load_qa_chain
from langchain_community.llms import HuggingFacePipeline  # Updated imports
import os
from wordcloud import WordCloud
import matplotlib.pyplot as plt
from transformers import BartForConditionalGeneration, BartTokenizer
from difflib import HtmlDiff, SequenceMatcher
import tempfile
from transformers import AutoModelForCausalLM, AutoTokenizer, AutoModelForTokenClassification
from transformers import pipeline
import nltk
import networkx as nx
from keybert import KeyBERT
from nltk.tokenize import word_tokenize


# Initialize KeyBERT for keyword extraction
kw_model = KeyBERT()


# Download NLTK punkt tokenizer for sentence splitting
nltk.download("punkt")
from nltk.tokenize import sent_tokenize


# <-----------------------------------Set up Streamlit app------------------------------------>
st.set_page_config(page_title="Corporate Training Knowledge Hub", layout="wide")
st.title("Corporate Training Knowledge Hub")

# <------------------------------------Initialize components------------------------------------->


# Initialize the LLaMA model using Ollama
llm = Ollama(model="llama3.2")  # Replace with your locally installed LLaMA model


# Load embedding model for document processing
embedding_model = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
vectorstore = None
document_store = []

# <---------------------------------------------Define tabs for functionalities------------------------------------>

tabs = st.tabs(["Upload Files", "Original Context", "Document Summarization", "Interactive Q&A", "Word Cloud", "Compare Docs", "Highlights","Course Path Generation"])

# <--------------------------------------------------Upload and process files------------------------------------->
def process_files(uploaded_files):
    global vectorstore
    for uploaded_file in uploaded_files:
        file_type = uploaded_file.name.split(".")[-1]
        combined_content = ""
        document = None

        with tempfile.NamedTemporaryFile(delete=False) as temp_file:
            temp_file.write(uploaded_file.read())
            temp_file_path = temp_file.name

        if file_type == "pdf":
            loader = PyPDFLoader(temp_file_path)
            all_pages = loader.load()
            combined_content = " ".join([page.page_content for page in all_pages])
            document = Document(page_content=combined_content, metadata={"name": uploaded_file.name})

        elif file_type == "pptx":
            presentation = Presentation(temp_file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        combined_content += shape.text_frame.text + " "
            document = Document(page_content=combined_content, metadata={"name": uploaded_file.name})

        elif file_type == "txt":
            with open(temp_file_path, "r", encoding="utf-8") as file:
                combined_content = file.read()
            document = Document(page_content=combined_content, metadata={"name": uploaded_file.name})

        elif file_type == "xlsx":
            excel_data = pd.read_excel(temp_file_path)
            combined_content = excel_data.to_string(index=False)
            document = Document(page_content=combined_content, metadata={"name": uploaded_file.name})

        if document is None:
            st.warning(f"File type '{file_type}' is not supported.")
            os.remove(temp_file_path)
            continue

        document_store.append(document)
        texts = [document.page_content]
        if vectorstore is None:
            vectorstore = FAISS.from_texts(texts, embedding_model)
        else:
            vectorstore.add_texts(texts)

        os.remove(temp_file_path)

# <----------------------------------------------------Summarization function------------------------------------->

def summarize_text_with_llama(text):
    """
    Summarizes the provided text using the locally running LLaMA 3.2 model.
    """
    # Prepare the context for summarization
    prompt = f"""
    Please summarize the following text:
    
    {text}
    """

    # Use the locally running LLaMA model to generate the summary
    response = llm(prompt)

    # Extract and return the generated summary
    summary = response.strip()
    return summary

# <------------------------------------------------------Interactive Q&A Functionality----------------------------------->

def answer_question_with_llama(question):
    """
    Answers a question using the locally installed LLaMA model with Ollama.
    """
    if not vectorstore:
        return "No documents indexed for retrieval. Please upload files first."

    # Retrieve relevant documents
    retriever = vectorstore.as_retriever()
    retrieved_docs = retriever.get_relevant_documents(question)

    if not retrieved_docs:
        return "No relevant documents found for your question."

    # Combine the content of retrieved documents
    combined_context = " ".join([doc.page_content for doc in retrieved_docs])

    # Prepare the prompt for the model
    prompt = f"Context: {combined_context}\n\nQuestion: {question}\nAnswer:"

    # Generate a response using the Ollama model
    response = llm(prompt)

    # Return the generated response
    return response

# <-------------------------------------------- Word Cloud Function---------------------------------->
def generate_word_cloud(text):
    """
    Generates a word cloud visualization for the provided text.
    """
    wordcloud = WordCloud(background_color="white", width=800, height=400, max_words=200, colormap="viridis").generate(text)
    plt.figure(figsize=(10, 5))
    plt.imshow(wordcloud, interpolation="bilinear")
    plt.axis("off")
    st.pyplot(plt)

# <------------------------------------------------------------Compare Documents------------------------------------------------------>
def compare_documents():
    """
    Compares the first two documents in the document store for similarities and differences.
    """
    if len(document_store) < 2:
        st.error("Please upload at least two documents to compare.")
        return

    doc1_content = document_store[0].page_content
    doc2_content = document_store[1].page_content

    differ = HtmlDiff()
    html_diff = differ.make_file(doc1_content.splitlines(), doc2_content.splitlines(), 
                                 fromdesc=document_store[0].metadata['name'], todesc=document_store[1].metadata['name'])
    st.write("### Comparison Result")
    st.components.v1.html(html_diff, height=600, scrolling=True)
    similarity = SequenceMatcher(None, doc1_content, doc2_content).ratio()
    st.write(f"### Similarity Score: {similarity:.2%}")


#< -----------------------------------------------------------Highlights-------------------------------------->


ner_model = pipeline("ner", model="dbmdz/bert-large-cased-finetuned-conll03-english", aggregation_strategy="simple")


def generate_description_with_ollama(entity_text, entity_type, context):
    """
    Generates concise descriptions for entities using Ollama.
    """
    prompt = f"""
    Entity: {entity_text} ({entity_type})
    Context: {context}
    Task: Provide a single-sentence refined description of this entity.
    """

    # Use the locally installed Ollama model
    response = llm(prompt)
    return response.strip()

def extract_highlights_with_ollama(text):
    """
    Extracts concise highlights and generates descriptions using Ollama 3.2.
    """
    entities = ner_model(text)
    valid_entity_types = {"PER", "ORG", "LOC", "GPE", "DATE"}
    seen_entities = set()
    highlights = []

    for entity in entities:
        entity_text = entity["word"]
        entity_type = entity["entity_group"]

        if entity_type in valid_entity_types and entity_text not in seen_entities:
            seen_entities.add(entity_text)

            # Extract context sentences (1-2 sentences only)
            context_sentences = [s for s in sent_tokenize(text) if entity_text in s]
            context = " ".join(context_sentences[:1]) if context_sentences else "No detailed context available."

            # Generate concise description
            refined_description = generate_description_with_ollama(entity_text, entity_type, context)
            highlights.append(f"{entity_text} ({entity_type}) - {refined_description}")

    return highlights




# <--------------------------------------------Course Path Generation Functions---------------------------------->
# def extract_unique_keywords(text, num_keywords=10):
#     """
#     Extracts unique course-related keywords using KeyBERT.
#     """
#     from keybert import KeyBERT  # Ensure KeyBERT is imported
#     kw_model = KeyBERT()

#     # Extract keywords and ensure they are unique
#     keywords = kw_model.extract_keywords(text, keyphrase_ngram_range=(1, 2), stop_words="english", top_n=num_keywords)
#     unique_keywords = list(set([kw[0] for kw in keywords]))  # Remove duplicates
#     return unique_keywords


# def extract_meaningful_keywords(text, num_keywords=10):
#     """
#     Extracts meaningful course-related keywords using LLaMA.
#     """
#     # Prepare the prompt for LLaMA to identify meaningful course-related keywords
#     prompt = f"""
#     Extract the most meaningful and structured course-related topics from the following text. 
#     Ensure the topics are actionable and related to well-defined subjects like "Machine Learning," "Python," "Java," etc.
#     Limit the output to {num_keywords} topics:
    
#     {text}
#     """

#     # Generate keywords using the locally running LLaMA model
#     response = llm(prompt)

#     # Extract the keywords from the response (assuming one topic per line or comma-separated)
#     keywords = [keyword.strip() for keyword in response.split("\n") if keyword.strip()]
#     return keywords[:num_keywords]  # Limit to the top num_keywords


# def parse_learning_path_to_graph(learning_path):
#     """
#     Parses the learning path text into a hierarchical graph structure for visualization.
#     Assumes LLaMA returns structured text with levels: Beginner, Intermediate, Advanced.
#     """
#     levels = {"Beginner": [], "Intermediate": [], "Advanced": []}
#     for line in learning_path.split("\n"):
#         if line.startswith("Beginner"):
#             levels["Beginner"] = [topic.strip() for topic in line[len("Beginner:"):].split(",") if topic.strip()]
#         elif line.startswith("Intermediate"):
#             levels["Intermediate"] = [topic.strip() for topic in line[len("Intermediate:"):].split(",") if topic.strip()]
#         elif line.startswith("Advanced"):
#             levels["Advanced"] = [topic.strip() for topic in line[len("Advanced:"):].split(",") if topic.strip()]
#     return levels

# def visualize_roadmap(keyword, levels):
#     """
#     Creates a roadmap-style visualization for the learning path using NetworkX.
#     """
#     graph = nx.DiGraph()

#     # Add nodes and edges for the hierarchical structure
#     for i, level in enumerate(["Beginner", "Intermediate", "Advanced"]):
#         if levels[level]:
#             for topic in levels[level]:
#                 graph.add_node(topic, level=level)
#                 if i > 0:  # Link nodes from the previous level
#                     for prev_topic in levels[["Beginner", "Intermediate", "Advanced"][i - 1]]:
#                         graph.add_edge(prev_topic, topic)

#     # Generate positions for the hierarchical layout
#     pos = nx.multipartite_layout(graph, subset_key="level")

#     # Draw the graph
#     plt.figure(figsize=(12, 8))
#     nx.draw(
#         graph, pos, with_labels=True, node_size=3000, node_color="skyblue", 
#         font_size=10, font_weight="bold", edge_color="gray", arrowsize=15
#     )
#     plt.title(f"Learning Path Roadmap for '{keyword}'")
#     st.pyplot(plt)

def extract_meaningful_keywords_with_llama(text, num_keywords=10):
    """
    Extracts structured, course-like keywords using LLaMA.
    Filters to ensure actionable topics are returned.
    """
    prompt = f"""
    From the following text, identify up to {num_keywords} concise, well-defined course topics.
    Only return short, actionable topics like "Python," "Machine Learning," "Deep Learning," etc.
    Avoid explanations, summaries, or full sentences.
    
    {text}
    """

    # Use the LLaMA model to generate the keywords
    response = llm(prompt)

    # Clean and extract keywords from the response
    keywords = [keyword.strip() for keyword in response.split(",") if len(keyword.strip()) > 0]
    return keywords[:num_keywords]  # Limit to the top num_keywords


def generate_learning_path_with_llama(keyword):
    """
    Generates a structured learning path for the given keyword using the LLaMA model.
    """
    prompt = f"""
    Create a structured learning roadmap for the topic '{keyword}' divided into three levels: Beginner, Intermediate, and Advanced.
    Each level should have 3-5 concise subtopics or key concepts to learn. Do not include explanations or long descriptions.
    Example format:
    Beginner: Topic 1, Topic 2
    Intermediate: Topic 3, Topic 4
    Advanced: Topic 5, Topic 6
    """

    # Use the LLaMA model to generate the learning path
    response = llm(prompt)

    # Return the structured learning path
    return response.strip()


def parse_and_clean_learning_path(learning_path):
    """
    Cleans and parses the learning path text into a hierarchical format for visualization.
    Handles irregular responses and ensures proper structure.
    """
    levels = {"Beginner": [], "Intermediate": [], "Advanced": []}

    # Parse the LLaMA response and map to levels
    for line in learning_path.split("\n"):
        if line.lower().startswith("beginner"):
            levels["Beginner"] = [topic.strip() for topic in line[len("Beginner:"):].split(",") if topic.strip()]
        elif line.lower().startswith("intermediate"):
            levels["Intermediate"] = [topic.strip() for topic in line[len("Intermediate:"):].split(",") if topic.strip()]
        elif line.lower().startswith("advanced"):
            levels["Advanced"] = [topic.strip() for topic in line[len("Advanced:"):].split(",") if topic.strip()]
    return levels

def visualize_roadmap_with_fallback(keyword, levels):
    """
    Creates a roadmap-style visualization or handles fallback in case of insufficient data.
    """
    graph = nx.DiGraph()

    # Add nodes and edges for hierarchical levels
    for i, level in enumerate(["Beginner", "Intermediate", "Advanced"]):
        if levels[level]:
            for topic in levels[level]:
                graph.add_node(topic, level=level)
                if i > 0:  # Link to the previous level
                    for prev_topic in levels[["Beginner", "Intermediate", "Advanced"][i - 1]]:
                        graph.add_edge(prev_topic, topic)

    if graph.number_of_nodes() == 0:
        st.warning(f"No structured roadmap found for '{keyword}'. Check LLaMA output for adjustments.")
        return

    # Generate hierarchical positions for the graph
    pos = nx.multipartite_layout(graph, subset_key="level")

    # Draw the graph
    plt.figure(figsize=(12, 8))
    nx.draw(
        graph, pos, with_labels=True, node_size=3000, node_color="lightblue",
        font_size=10, font_weight="bold", edge_color="gray", arrowsize=15
    )
    plt.title(f"Learning Path Roadmap for '{keyword}'")
    st.pyplot(plt)


# <-------------------------------------------------------Main App-------------------------------->
st.sidebar.header("Welcome!")
st.sidebar.info("Upload corporate training documents, explore their contents, get concise summaries, generate word clouds, and ask interactive questions!")

with tabs[0]:
    st.header("Upload Files")
    uploaded_files = st.file_uploader("Upload corporate documents (PDF, PPTX, TXT, XLSX)", 
                                      type=["pdf", "pptx", "txt", "xlsx"], accept_multiple_files=True)
    if uploaded_files:
        process_files(uploaded_files)
        st.success("Files processed successfully!")

with tabs[1]:
    st.header("Original Context")
    if document_store:
        for doc in document_store:
            st.write(f"### {doc.metadata['name']}")
            st.text_area(f"Content_{doc.metadata['name']}", doc.page_content, height=300)
    else:
        st.info("Please upload files to display their content.")

with tabs[2]:
    st.header("Document Summarization")
    if document_store:
        for doc in document_store:
            st.write(f"### {doc.metadata['name']}")
            summary = summarize_text_with_llama(doc.page_content)
            st.write("### Summary:")
            st.write(summary)
    else:
        st.info("Please upload files to summarize.")

with tabs[3]:
    st.header("Interactive Q&A")
    question = st.text_input("Ask a question about the uploaded documents:")
    if question:
        answer = answer_question_with_llama(question)
        st.write("### Answer:")
        st.write(answer)

with tabs[4]:
    st.header("Word Cloud")
    if document_store:
        text_data = " ".join([doc.page_content for doc in document_store])
        st.write("### Word Cloud")
        generate_word_cloud(text_data)
    else:
        st.info("Please upload files to generate a word cloud.")

with tabs[5]:
    st.header("Compare Documents")
    compare_documents()

with tabs[6]:
    st.header("Highlights (Concise Contextual Insights)")

    if document_store:
        for doc in document_store:
            st.subheader(f"Document: {doc.metadata['name']}")
            highlights = extract_highlights_with_ollama(doc.page_content)
            if highlights:
                for i, highlight in enumerate(highlights, start=1):
                    st.markdown(f"**{i}. {highlight}**")
            else:
                st.write("No significant entities or highlights found.")
    else:
        st.info("No documents uploaded yet.")


with tabs[7]:
    st.header("Learning Path Roadmap")
    if document_store:
        for doc in document_store:
            st.subheader(f"Document: {doc.metadata['name']}")

            # Extract meaningful course-like keywords
            keywords = extract_meaningful_keywords_with_llama(doc.page_content)
            st.write("### Keywords Identified:")
            st.write(", ".join(keywords))

            # Generate and visualize roadmaps for each keyword
            for keyword in keywords:
                st.write(f"### Roadmap for {keyword}:")
                learning_path = generate_learning_path_with_llama(keyword)  # Get LLaMA response
                levels = parse_and_clean_learning_path(learning_path)  # Parse response into levels
                visualize_roadmap_with_fallback(keyword, levels)  # Visualize roadmap
    else:
        st.info("Please upload documents to generate learning paths.")
